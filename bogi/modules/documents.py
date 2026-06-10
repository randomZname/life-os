"""Document ingestion, chunking, embedding, retrieval (pgvector)."""

from __future__ import annotations

import hashlib
import logging
from functools import lru_cache
from pathlib import Path
from typing import Any

from langchain_text_splitters import RecursiveCharacterTextSplitter
from sqlalchemy import func, select

from bogi.config import settings
from bogi.db import get_session
from bogi.models import Chunk, Course, Document

logger = logging.getLogger(__name__)


# ---------- Embedding model (lazy singleton) ----------


@lru_cache(maxsize=1)
def _get_embedder():
    """Lazy-load sentence-transformers модел. Тежък import — само при първо ползване."""
    from sentence_transformers import SentenceTransformer

    logger.info("Loading embedding model: %s", settings.embedding_model)
    model = SentenceTransformer(settings.embedding_model)
    actual_dim = model.get_sentence_embedding_dimension()
    if actual_dim != settings.embedding_dimension:
        logger.warning(
            "Embedding dimension mismatch: model=%d settings=%d. Update EMBEDDING_DIMENSION.",
            actual_dim,
            settings.embedding_dimension,
        )
    return model


def embed_texts(texts: list[str]) -> list[list[float]]:
    """Encode list of texts -> list of embedding vectors."""
    model = _get_embedder()
    embeddings = model.encode(texts, convert_to_numpy=True, normalize_embeddings=True)
    return embeddings.tolist()


# ---------- Text extraction ----------


def _extract_text(file_path: Path) -> str:
    """Извличане на текст от PDF, DOCX, PPTX, TXT, MD."""
    suffix = file_path.suffix.lower()
    try:
        if suffix == ".pdf":
            import pymupdf

            doc = pymupdf.open(str(file_path))
            text = "\n\n".join(page.get_text() for page in doc)
            doc.close()
            return text
        if suffix in {".docx"}:
            from docx import Document as DocxDocument

            d = DocxDocument(str(file_path))
            return "\n".join(p.text for p in d.paragraphs)
        if suffix in {".pptx"}:
            from pptx import Presentation

            prs = Presentation(str(file_path))
            parts: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
            return "\n".join(parts)
        if suffix in {".txt", ".md"}:
            return file_path.read_text(encoding="utf-8", errors="ignore")
    except Exception as exc:
        logger.exception("Text extraction failed for %s: %s", file_path, exc)
        return ""

    logger.warning("Unsupported file type: %s", suffix)
    return ""


def _file_hash(file_path: Path) -> str:
    h = hashlib.sha256()
    with file_path.open("rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


# ---------- Ingestion ----------


async def document_ingest(
    file_path: str,
    source: str = "manual",
    course_id: int | None = None,
) -> dict[str, Any]:
    """Chunk + embed + insert. Идемпотентно по content_hash."""
    path = Path(file_path).resolve()
    if not path.exists():
        return {"ok": False, "error": f"Файлът не съществува: {file_path}"}

    text = _extract_text(path)
    if not text.strip():
        return {"ok": False, "error": "Не успях да извлека текст от файла"}

    sha = _file_hash(path)

    async with get_session() as session:
        # Idempotency: ако вече е ingest-нат, връщаме no-op
        existing = await session.execute(
            select(Document).where(Document.content_hash == sha)
        )
        existing_doc = existing.scalar_one_or_none()
        if existing_doc:
            return {
                "ok": True,
                "noop": True,
                "document_id": existing_doc.id,
                "message": "Документът вече е ingest-нат (същият content_hash).",
            }

        doc = Document(
            course_id=course_id,
            file_path=str(path),
            title=path.stem,
            source=source,
            content_hash=sha,
            extracted_text=text[:1_000_000],  # cap при много дълъг текст
        )
        session.add(doc)
        await session.flush()

        # Chunking
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            separators=["\n\n", "\n", ". ", " ", ""],
        )
        chunks_text = splitter.split_text(text)
        if not chunks_text:
            return {"ok": True, "document_id": doc.id, "chunks": 0}

        embeddings = embed_texts(chunks_text)

        for idx, (txt, emb) in enumerate(zip(chunks_text, embeddings, strict=True)):
            session.add(
                Chunk(
                    document_id=doc.id,
                    chunk_idx=idx,
                    text=txt,
                    embedding=emb,
                )
            )

        return {
            "ok": True,
            "document_id": doc.id,
            "chunks": len(chunks_text),
            "title": doc.title,
        }


# ---------- Retrieval ----------


# pg_trgm WORD-similarity floor for the trigram candidate query (0..1). We use
# `word_similarity(query, text)` (best match of the query inside the longer chunk)
# — NOT plain `similarity()`, which compares the whole query against the whole
# chunk and is ~0 for a short query vs a long chunk (verified live: an exact
# "Бойс-Код нормална форма" match scored similarity=0.14 but word_similarity=1.0).
# 0.5 cleanly separates real lexical hits from trigram noise on unrelated chunks.
TRGM_SIM_THRESHOLD = 0.5


def _merge_hybrid(
    vector_rows: list[dict[str, Any]],
    trigram_rows: list[dict[str, Any]],
    k: int,
) -> list[dict[str, Any]]:
    """Union vector + trigram hits by ``chunk_id``, keep the higher score.

    Pure helper (no DB / no I/O) so it is unit-testable in isolation:
    - dedup by ``chunk_id``: a chunk found by both paths appears once;
    - its score is ``max(vector_score, trigram_score)``;
    - sort by ``score`` descending; return the top-``k``.
    The dict shape returned is identical to the per-row shape produced by the
    vector and trigram queries (``agent.py`` depends on these exact keys).
    """
    merged: dict[Any, dict[str, Any]] = {}
    for row in [*vector_rows, *trigram_rows]:
        cid = row["chunk_id"]
        existing = merged.get(cid)
        if existing is None or row["score"] > existing["score"]:
            # Keep the row carrying the higher score (so `score` is the max).
            merged[cid] = row
    ordered = sorted(merged.values(), key=lambda r: r["score"], reverse=True)
    return ordered[:k]


async def document_search(query: str, k: int = 5) -> list[dict[str, Any]]:
    """Hybrid search: pgvector (cosine) ∪ pg_trgm keyword. Top-k chunks + метаданни.

    Vector half ranks by embedding cosine similarity; trigram half ranks by
    ``func.word_similarity(query, Chunk.text)`` (finds the query as a word-set
    inside the longer chunk; NOT plain ``similarity()`` which is ~0 for a short
    query vs a long chunk, and NOT the bare ``%`` operator — it silently no-ops
    under asyncpg; see D-019 and the TRGM_SIM_THRESHOLD note). Results are
    merged by ``_merge_hybrid`` (union by chunk_id, max score, top-k). If the
    trigram query errors (extension/index missing), we log and fall back to
    vector-only so search never breaks. Returned dict keys are unchanged:
    chunk_id, document_id, title, file_path, chunk_idx, text, score.
    """
    [query_emb] = embed_texts([query])

    async with get_session() as session:
        # --- Vector candidates: cosine distance (<=> в pgvector) ---
        vec_stmt = (
            select(
                Chunk.id,
                Chunk.text,
                Chunk.chunk_idx,
                Document.id.label("document_id"),
                Document.title,
                Document.file_path,
                Chunk.embedding.cosine_distance(query_emb).label("distance"),
            )
            .join(Document, Chunk.document_id == Document.id)
            .order_by(Chunk.embedding.cosine_distance(query_emb))
            .limit(k)
        )
        vec_result = await session.execute(vec_stmt)
        vector_rows = [
            {
                "chunk_id": r.id,
                "document_id": r.document_id,
                "title": r.title,
                "file_path": r.file_path,
                "chunk_idx": r.chunk_idx,
                "text": r.text,
                "score": 1.0 - float(r.distance),
            }
            for r in vec_result.all()
        ]

        # --- Trigram candidates: pg_trgm WORD similarity (uses ix_chunks_text_trgm).
        # word_similarity(query, text) = best match of the query within the longer
        # chunk. NOT plain similarity() (≈0 for short-query/long-chunk) and NOT the
        # bare `%` operator (renders as unknown `%%` under asyncpg → silent no-op,
        # decision D-019).
        sim = func.word_similarity(query, Chunk.text)
        trg_stmt = (
            select(
                Chunk.id,
                Chunk.text,
                Chunk.chunk_idx,
                Document.id.label("document_id"),
                Document.title,
                Document.file_path,
                sim.label("similarity"),
            )
            .join(Document, Chunk.document_id == Document.id)
            .where(sim > TRGM_SIM_THRESHOLD)
            .order_by(sim.desc())
            .limit(k)
        )
        try:
            trg_result = await session.execute(trg_stmt)
            trigram_rows = [
                {
                    "chunk_id": r.id,
                    "document_id": r.document_id,
                    "title": r.title,
                    "file_path": r.file_path,
                    "chunk_idx": r.chunk_idx,
                    "text": r.text,
                    "score": float(r.similarity),
                }
                for r in trg_result.all()
            ]
        except Exception:
            # Trigram search failed (extension/index missing?) — fall back to
            # vector-only. Search must never break.
            logger.warning("hybrid trigram candidate query failed", exc_info=True)
            trigram_rows = []

    return _merge_hybrid(vector_rows, trigram_rows, k)


async def document_list(limit: int = 50) -> list[dict[str, Any]]:
    """Списък на ingest-натите документи."""
    async with get_session() as session:
        stmt = (
            select(Document, Course.name.label("course_name"))
            .outerjoin(Course, Document.course_id == Course.id)
            .order_by(Document.created_at.desc())
            .limit(limit)
        )
        result = await session.execute(stmt)
        return [
            {
                "id": doc.id,
                "title": doc.title,
                "file_path": doc.file_path,
                "source": doc.source,
                "course": course_name,
                "created_at": doc.created_at.isoformat(),
            }
            for doc, course_name in result.all()
        ]


async def document_read(doc_id: int, max_chars: int = 50_000) -> dict[str, Any]:
    """Връща пълния (truncated) текст на документ."""
    async with get_session() as session:
        doc = await session.get(Document, doc_id)
        if not doc:
            return {"ok": False, "error": "Не е намерен документ"}
        text = doc.extracted_text or ""
        return {
            "ok": True,
            "id": doc.id,
            "title": doc.title,
            "text": text[:max_chars],
            "truncated": len(text) > max_chars,
        }


# ---------- Helpers ----------


async def upsert_course(fmi_id: str, name: str, url: str | None = None) -> int:
    """Insert-or-update course по fmi_id. Връща course.id."""
    async with get_session() as session:
        existing = await session.execute(select(Course).where(Course.fmi_id == fmi_id))
        course = existing.scalar_one_or_none()
        if course:
            course.name = name
            if url:
                course.url = url
            return course.id
        course = Course(fmi_id=fmi_id, name=name, url=url)
        session.add(course)
        await session.flush()
        return course.id
